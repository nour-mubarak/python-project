#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
LinguaEval — Report Generator
================================

Generates a professional Multilingual AI Readiness Report (DOCX or PDF)
from evaluation results JSON.

Usage:
    python generate_report.py --results results/client_results.json
    python generate_report.py --results results/client_results.json --output reports/client_report.docx
    python generate_report.py --results results/client_results.json --format pdf
"""

import argparse
import json
import os
import sys
from datetime import datetime
from typing import Dict, List, Optional

from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.section import WD_ORIENT
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml

# PDF generation
from fpdf import FPDF

# Slide deck generation
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ═══════════════════════════════════════════════════════════════
# DESIGN CONSTANTS
# ═══════════════════════════════════════════════════════════════

NAVY = RGBColor(0x0C, 0x2D, 0x48)
ACCENT = RGBColor(0x02, 0x84, 0xC7)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x16, 0x65, 0x34)
RED = RGBColor(0xB9, 0x1C, 0x1C)
ORANGE = RGBColor(0xC2, 0x41, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = "EBF5FB"

FONT_BODY = "Calibri"
FONT_SIZE = Pt(10.5)


def severity_color(severity: str) -> RGBColor:
    """Get colour for severity level."""
    return {
        "low": GREEN,
        "medium": ORANGE,
        "high": RED,
        "critical": RED,
    }.get(severity, GRAY)


def score_to_status(score: float) -> str:
    """Convert average score to deployment status."""
    if score >= 80:
        return "Ready for Pilot"
    elif score >= 65:
        return "Restricted Pilot Only"
    return "Not Ready"


def score_to_color(score: float) -> RGBColor:
    if score >= 80:
        return GREEN
    elif score >= 65:
        return ORANGE
    return RED


# ═══════════════════════════════════════════════════════════════
# DOCUMENT HELPERS
# ═══════════════════════════════════════════════════════════════


def set_cell_shading(cell, color_hex: str):
    """Set background colour on a table cell."""
    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}"/>')
    cell._tc.get_or_add_tcPr().append(shading)


def add_styled_paragraph(
    doc,
    text,
    style="Normal",
    bold=False,
    color=None,
    size=None,
    alignment=None,
    space_after=None,
    space_before=None,
):
    """Add a paragraph with custom styling."""
    p = doc.add_paragraph()
    if alignment:
        p.alignment = alignment
    if space_after is not None:
        p.paragraph_format.space_after = Pt(space_after)
    if space_before is not None:
        p.paragraph_format.space_before = Pt(space_before)

    run = p.add_run(text)
    run.font.name = FONT_BODY
    run.font.size = size or FONT_SIZE
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return p


def add_rich_paragraph(doc, runs_data, alignment=None, space_after=None):
    """Add a paragraph with mixed formatting."""
    p = doc.add_paragraph()
    if alignment:
        p.alignment = alignment
    if space_after is not None:
        p.paragraph_format.space_after = Pt(space_after)

    for rd in runs_data:
        run = p.add_run(rd.get("text", ""))
        run.font.name = FONT_BODY
        run.font.size = rd.get("size", FONT_SIZE)
        run.font.bold = rd.get("bold", False)
        run.font.italic = rd.get("italic", False)
        if rd.get("color"):
            run.font.color.rgb = rd["color"]
    return p


def add_key_value_table(doc, data: List[tuple], col_widths=(Inches(2), Inches(4.5))):
    """Add a simple two-column key-value table."""
    table = doc.add_table(rows=len(data), cols=2)
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    for i, (key, value) in enumerate(data):
        # Key cell
        cell_k = table.cell(i, 0)
        cell_k.width = col_widths[0]
        p = cell_k.paragraphs[0]
        run = p.add_run(key)
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = NAVY
        set_cell_shading(cell_k, "F8FAFC")

        # Value cell
        cell_v = table.cell(i, 1)
        cell_v.width = col_widths[1]
        p = cell_v.paragraphs[0]
        run = p.add_run(str(value))
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.color.rgb = DARK

    return table


def add_scores_table(doc, models_data: dict, dimension_names: list):
    """Add a model comparison scores table."""
    models = list(models_data.keys())
    n_cols = (
        1 + len(models) * 2 + len(models)
    )  # dim + (en + ar) per model + gap per model

    # Simplified: dim | model1_en | model1_ar | model1_gap | model2_en | ...
    header_cols = ["Dimension"]
    for m in models:
        short_name = m.split("-")[0].upper() if "-" in m else m.upper()
        header_cols.extend(
            [f"{short_name} EN", f"{short_name} AR", f"{short_name} Gap"]
        )

    table = doc.add_table(rows=1 + len(dimension_names), cols=len(header_cols))
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    # Header row
    for j, header_text in enumerate(header_cols):
        cell = table.cell(0, j)
        p = cell.paragraphs[0]
        run = p.add_run(header_text)
        run.font.name = FONT_BODY
        run.font.size = Pt(8.5)
        run.font.bold = True
        run.font.color.rgb = WHITE
        set_cell_shading(cell, "0C2D48")

    # Data rows
    for i, dim in enumerate(dimension_names):
        row_idx = i + 1
        # Dimension name
        cell = table.cell(row_idx, 0)
        p = cell.paragraphs[0]
        run = p.add_run(dim.replace("_", " ").title())
        run.font.name = FONT_BODY
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = DARK

        if i % 2 == 0:
            set_cell_shading(cell, "F8FAFC")

        col_offset = 1
        for m in models:
            m_data = models_data.get(m, {})

            # Consistency is a pairwise metric (compares EN vs AR) - handle specially
            if dim == "consistency":
                # Get the consistency score (stored under ar since it's cross-lingual)
                consistency_score = (
                    m_data.get("ar", {}).get(dim, {}).get("average", "N/A")
                )

                # EN column: show "Pairwise" indicator
                cell_en = table.cell(row_idx, col_offset)
                p_en = cell_en.paragraphs[0]
                p_en.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = p_en.add_run("Pairwise")
                run.font.name = FONT_BODY
                run.font.size = Pt(8)
                run.font.italic = True
                run.font.color.rgb = GRAY
                if i % 2 == 0:
                    set_cell_shading(cell_en, "F8FAFC")

                # AR column: show the pairwise consistency score
                cell_ar = table.cell(row_idx, col_offset + 1)
                p_ar = cell_ar.paragraphs[0]
                p_ar.alignment = WD_ALIGN_PARAGRAPH.CENTER
                if isinstance(consistency_score, (int, float)):
                    run = p_ar.add_run(f"{consistency_score:.1f}%")
                    run.font.name = FONT_BODY
                    run.font.size = Pt(9)
                    run.font.bold = True
                    run.font.color.rgb = score_to_color(consistency_score)
                else:
                    run = p_ar.add_run("N/A")
                    run.font.name = FONT_BODY
                    run.font.size = Pt(9)
                    run.font.color.rgb = GRAY
                if i % 2 == 0:
                    set_cell_shading(cell_ar, "F8FAFC")

                # Gap column: show dash (consistency IS the cross-lingual score)
                cell_gap = table.cell(row_idx, col_offset + 2)
                p_gap = cell_gap.paragraphs[0]
                p_gap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = p_gap.add_run("—")
                run.font.name = FONT_BODY
                run.font.size = Pt(9)
                run.font.color.rgb = GRAY
                if i % 2 == 0:
                    set_cell_shading(cell_gap, "F8FAFC")
            else:
                # Standard dimension: show EN, AR, and Gap
                en_score = m_data.get("en", {}).get(dim, {}).get("average", "N/A")
                ar_score = m_data.get("ar", {}).get(dim, {}).get("average", "N/A")
                gap_data = m_data.get("cross_lingual_gap", {}).get(dim, {})
                gap = gap_data.get("gap", "N/A")

                for k, (val, is_gap) in enumerate(
                    [(en_score, False), (ar_score, False), (gap, True)]
                ):
                    cell = table.cell(row_idx, col_offset + k)
                    p = cell.paragraphs[0]
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

                    if isinstance(val, (int, float)):
                        text = f"{val:.1f}%"
                        run = p.add_run(text)
                        run.font.name = FONT_BODY
                        run.font.size = Pt(9)
                        run.font.bold = True
                        if is_gap:
                            sev = gap_data.get("severity", "low")
                            run.font.color.rgb = severity_color(sev)
                        else:
                            run.font.color.rgb = score_to_color(val)
                    else:
                        run = p.add_run("N/A")
                        run.font.name = FONT_BODY
                        run.font.size = Pt(9)
                        run.font.color.rgb = GRAY

                    if i % 2 == 0:
                        set_cell_shading(cell, "F8FAFC")

            col_offset += 3

    return table


# ═══════════════════════════════════════════════════════════════
# REPORT GENERATOR
# ═══════════════════════════════════════════════════════════════


class ReportGenerator:
    """Generates the Multilingual AI Readiness Report from evaluation results."""

    def __init__(self, results_path: str):
        """Load results from JSON file."""
        with open(results_path, encoding="utf-8") as f:
            self.data = json.load(f)

        self.metadata = self.data.get("metadata", {})
        self.aggregates = self.data.get("aggregates", {})
        self.detailed = self.data.get("detailed_results", [])

        self.client_name = self.metadata.get("client_name", "Client")
        self.sector = self.metadata.get("sector", "general")
        self.models = self.metadata.get("models", [])
        self.timestamp = self.metadata.get("timestamp", datetime.now().isoformat())

    def _calculate_overall_scores(self) -> dict:
        """Calculate overall scores per model."""
        overall = {}
        for model_id, data in self.aggregates.items():
            en_scores = [v["average"] for v in data.get("en", {}).values()]
            ar_scores = [v["average"] for v in data.get("ar", {}).values()]
            gaps = [v["gap"] for v in data.get("cross_lingual_gap", {}).values()]

            en_avg = sum(en_scores) / len(en_scores) if en_scores else 0
            ar_avg = sum(ar_scores) / len(ar_scores) if ar_scores else 0
            avg_gap = sum(gaps) / len(gaps) if gaps else 0
            combined = (en_avg + ar_avg) / 2

            # Find critical gaps (>15%) for risk assessment
            critical_gaps = [
                (dim, gd["gap"])
                for dim, gd in data.get("cross_lingual_gap", {}).items()
                if gd.get("severity") in ["critical", "high"]
            ]

            overall[model_id] = {
                "en_average": round(en_avg, 1),
                "ar_average": round(ar_avg, 1),
                "combined_average": round(combined, 1),
                "avg_gap": round(avg_gap, 1),
                "status": score_to_status(combined),
                "critical_gaps": critical_gaps,
            }
        return overall

    def _get_code_switching_findings(self) -> List[dict]:
        """Extract code-switching issues from fluency flags."""
        findings = []
        for result in self.detailed:
            for mid, model_data in result.get("model_scores", {}).items():
                ar_data = model_data.get("ar", {})
                for score_entry in ar_data.get("scores", []):
                    if score_entry["dimension"] == "fluency":
                        for flag in score_entry.get("flags", []):
                            if (
                                "code-switching" in flag.lower()
                                or "latin" in flag.lower()
                            ):
                                findings.append(
                                    {
                                        "model": mid,
                                        "prompt_id": result["prompt_id"],
                                        "flag": flag,
                                        "response_preview": ar_data.get("response", "")[
                                            :150
                                        ],
                                    }
                                )
        return findings

    def _get_recommended_model(self, overall: dict) -> tuple:
        """Determine which model to recommend and why."""
        if not overall:
            return None, ""

        # Score models on multiple factors
        model_scores = {}
        for model_id, scores in overall.items():
            # Lower gap is better, higher combined is better
            # Critical gaps are weighted heavily negative
            score = (
                scores["combined_average"]
                - (scores["avg_gap"] * 2)
                - (len(scores.get("critical_gaps", [])) * 10)
            )
            model_scores[model_id] = score

        best_model = max(model_scores.items(), key=lambda x: x[1])
        best_data = overall[best_model[0]]

        # Generate reason
        reasons = []
        if best_data["avg_gap"] < 5:
            reasons.append("strong bilingual parity")
        if not best_data.get("critical_gaps"):
            reasons.append("no critical cross-lingual gaps")
        if best_data["ar_average"] > 60:
            reasons.append("acceptable Arabic performance")

        reason = ", ".join(reasons) if reasons else "highest overall score"
        return best_model[0], reason

    def _collect_flags(
        self, model_id: str = None, dimension: str = None, language: str = None
    ) -> List[dict]:
        """Collect all flags from detailed results, optionally filtered."""
        flags = []
        for result in self.detailed:
            for mid, model_data in result.get("model_scores", {}).items():
                if model_id and mid != model_id:
                    continue
                for lang, lang_data in model_data.items():
                    if language and lang != language:
                        continue
                    for score_entry in lang_data.get("scores", []):
                        if dimension and score_entry["dimension"] != dimension:
                            continue
                        for flag in score_entry.get("flags", []):
                            flags.append(
                                {
                                    "prompt_id": result["prompt_id"],
                                    "model": mid,
                                    "language": lang,
                                    "dimension": score_entry["dimension"],
                                    "severity": score_entry["severity"],
                                    "flag": flag,
                                    "prompt_en": result.get("prompt_en", ""),
                                    "prompt_ar": result.get("prompt_ar", ""),
                                    "response": lang_data.get("response", "")[:200],
                                }
                            )
        return flags

    def generate(self, output_path: str):
        """Generate the complete Readiness Report."""
        doc = Document()

        # ── Page setup ──
        section = doc.sections[0]
        section.page_width = Cm(21)
        section.page_height = Cm(29.7)
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

        # ── Default styles ──
        style = doc.styles["Normal"]
        style.font.name = FONT_BODY
        style.font.size = FONT_SIZE
        style.font.color.rgb = DARK
        style.paragraph_format.space_after = Pt(6)
        style.paragraph_format.line_spacing = 1.15

        # Configure heading styles
        for level, (size, color) in enumerate(
            [(Pt(22), NAVY), (Pt(16), ACCENT), (Pt(13), DARK)], 1
        ):
            heading_style = doc.styles[f"Heading {level}"]
            heading_style.font.name = FONT_BODY
            heading_style.font.size = size
            heading_style.font.color.rgb = color
            heading_style.font.bold = True
            heading_style.paragraph_format.space_before = Pt(18 if level == 1 else 14)
            heading_style.paragraph_format.space_after = Pt(8)

        overall = self._calculate_overall_scores()

        # ═══════════════════════════════════════
        # COVER PAGE
        # ═══════════════════════════════════════
        for _ in range(6):
            doc.add_paragraph()

        add_styled_paragraph(
            doc,
            "MULTILINGUAL AI",
            bold=True,
            color=NAVY,
            size=Pt(28),
            alignment=WD_ALIGN_PARAGRAPH.CENTER,
            space_after=0,
        )
        add_styled_paragraph(
            doc,
            "READINESS REPORT",
            bold=True,
            color=NAVY,
            size=Pt(28),
            alignment=WD_ALIGN_PARAGRAPH.CENTER,
            space_after=8,
        )

        # Subtitle with use case context
        use_case = self.metadata.get("use_case", "General")
        add_styled_paragraph(
            doc,
            f"Arabic-English Model Evaluation for {use_case}",
            color=GRAY,
            size=Pt(12),
            alignment=WD_ALIGN_PARAGRAPH.CENTER,
            space_after=16,
        )

        add_styled_paragraph(
            doc,
            f"Prepared for {self.client_name}",
            color=ACCENT,
            size=Pt(14),
            alignment=WD_ALIGN_PARAGRAPH.CENTER,
            bold=True,
            space_after=24,
        )

        # Divider
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run("─" * 40)
        run.font.color.rgb = ACCENT
        run.font.size = Pt(10)

        date_str = datetime.fromisoformat(self.timestamp).strftime("%B %Y")
        add_styled_paragraph(
            doc,
            date_str,
            color=GRAY,
            size=Pt(12),
            alignment=WD_ALIGN_PARAGRAPH.CENTER,
            space_after=6,
        )
        add_styled_paragraph(
            doc,
            "LinguaEval",
            color=GRAY,
            size=Pt(12),
            alignment=WD_ALIGN_PARAGRAPH.CENTER,
            space_after=4,
        )
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run("Confidential")
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.color.rgb = GRAY
        run.font.italic = True

        doc.add_page_break()

        # ═══════════════════════════════════════
        # EXECUTIVE SUMMARY
        # ═══════════════════════════════════════
        doc.add_heading("1. Executive Summary", level=1)

        doc.add_paragraph(
            f"This report presents the results of a Multilingual AI Readiness Assessment "
            f"conducted for {self.client_name}. The evaluation tested {len(self.models)} AI model(s) "
            f"across Arabic and English on {self.metadata.get('total_prompts', 'N/A')} prompts "
            f"covering {self.sector} use cases."
        )

        # Overall status per model
        doc.add_heading("Deployment Readiness Summary", level=2)

        status_data = []
        for model_id, scores in overall.items():
            short = (
                model_id.split("-")[0].upper() if "-" in model_id else model_id.upper()
            )
            status_data.append(
                (
                    short,
                    f"{scores['en_average']}%",
                    f"{scores['ar_average']}%",
                    f"{scores['avg_gap']}%",
                    scores["status"],
                )
            )

        table = doc.add_table(rows=1 + len(status_data), cols=5)
        table.style = "Table Grid"
        table.alignment = WD_TABLE_ALIGNMENT.CENTER

        headers = [
            "Model",
            "English Score",
            "Arabic Score",
            "Avg Gap",
            "Readiness Status",
        ]
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            p = cell.paragraphs[0]
            run = p.add_run(h)
            run.font.name = FONT_BODY
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = WHITE
            set_cell_shading(cell, "0C2D48")

        for i, row_data in enumerate(status_data):
            for j, val in enumerate(row_data):
                cell = table.cell(i + 1, j)
                p = cell.paragraphs[0]
                run = p.add_run(str(val))
                run.font.name = FONT_BODY
                run.font.size = Pt(9.5)
                if j == 4:  # Status column
                    run.font.bold = True
                    if "Ready for Pilot" == val:
                        run.font.color.rgb = GREEN
                    elif "Restricted" in val:
                        run.font.color.rgb = ORANGE
                    else:
                        run.font.color.rgb = RED
                elif j in (1, 2):
                    score_val = float(val.replace("%", ""))
                    run.font.color.rgb = score_to_color(score_val)
                    run.font.bold = True

        doc.add_paragraph()

        # Key finding - use new recommendation logic
        recommended_model, reason = self._get_recommended_model(overall)
        if recommended_model:
            short_name = (
                recommended_model.split("-")[0].upper()
                if "-" in recommended_model
                else recommended_model.upper()
            )
            doc.add_heading("Recommended Model", level=2)
            p = doc.add_paragraph()
            run = p.add_run(f"{short_name}")
            run.font.bold = True
            run.font.color.rgb = GREEN
            run.font.size = Pt(12)
            run = p.add_run(f" — {reason}")
            run.font.size = Pt(10.5)

        # Critical findings - headline issues
        doc.add_heading("Critical Findings", level=2)

        # Check for code-switching
        code_switching = self._get_code_switching_findings()
        if code_switching:
            models_with_cs = set(f["model"] for f in code_switching)
            p = doc.add_paragraph()
            run = p.add_run("⚠ Code-Switching Detected: ")
            run.font.bold = True
            run.font.color.rgb = RED
            run.font.size = Pt(10)
            run = p.add_run(
                f"{len(code_switching)} instance(s) where Arabic prompts received responses with "
                f"significant English/Latin text. Affected model(s): {', '.join(models_with_cs)}. "
                f"This is a material Arabic-language quality risk requiring validation controls."
            )
            run.font.size = Pt(10)

        # Check for critical cross-lingual gaps
        critical_findings = []
        for model_id, scores in overall.items():
            for dim, gap in scores.get("critical_gaps", []):
                short = (
                    model_id.split("-")[0].upper()
                    if "-" in model_id
                    else model_id.upper()
                )
                critical_findings.append(f"{short}: {dim} gap of {gap:.1f}%")

        if critical_findings:
            p = doc.add_paragraph()
            run = p.add_run("⚠ Critical Cross-Lingual Gaps: ")
            run.font.bold = True
            run.font.color.rgb = RED
            run.font.size = Pt(10)
            run = p.add_run("; ".join(critical_findings))
            run.font.size = Pt(10)

        if not code_switching and not critical_findings:
            doc.add_paragraph("No critical findings requiring immediate attention.")

        doc.add_page_break()

        # ═══════════════════════════════════════
        # METHODOLOGY
        # ═══════════════════════════════════════
        doc.add_heading("2. Methodology", level=1)

        doc.add_paragraph(
            "This assessment used the LinguaEval evaluation framework, which tests AI models "
            "across six dimensions in both Arabic and English using sector-specific prompt packs."
        )

        doc.add_heading("Evaluation Parameters", level=2)
        add_key_value_table(
            doc,
            [
                ("Client", self.client_name),
                ("Sector", self.sector.title()),
                ("Use Case", self.metadata.get("use_case", "General")),
                ("Models Evaluated", ", ".join(self.models)),
                ("Languages", "English, Arabic"),
                ("Total Prompts", str(self.metadata.get("total_prompts", "N/A"))),
                (
                    "Evaluation Date",
                    datetime.fromisoformat(self.timestamp).strftime("%d %B %Y"),
                ),
                (
                    "Privacy Mode",
                    self.metadata.get("privacy_mode", "mode_a")
                    .replace("_", " ")
                    .title(),
                ),
            ],
        )

        doc.add_paragraph()

        doc.add_heading("Evaluation Dimensions", level=2)
        dimensions_desc = [
            (
                "Factual Accuracy",
                "Verifies responses against known ground truth and key facts.",
            ),
            (
                "Gender Bias",
                "Detects gendered defaults, stereotypical language, and bias asymmetries across languages.",
            ),
            (
                "Hallucination",
                "Identifies fabricated facts, unverified claims, and contradictions.",
            ),
            (
                "Cross-Lingual Consistency",
                "Measures whether Arabic and English responses provide equivalent information.",
            ),
            (
                "Cultural Sensitivity",
                "Evaluates cultural appropriateness and avoidance of stereotypes.",
            ),
            (
                "Fluency & Coherence",
                "Assesses language quality, completeness, and code-switching.",
            ),
        ]

        for dim_name, dim_desc in dimensions_desc:
            p = doc.add_paragraph()
            run = p.add_run(f"{dim_name}: ")
            run.font.bold = True
            run.font.color.rgb = NAVY
            run.font.size = Pt(10)
            run = p.add_run(dim_desc)
            run.font.size = Pt(10)

        doc.add_heading("Scoring Thresholds", level=2)
        doc.add_paragraph(
            "Within the LinguaEval framework, a cross-lingual gap above 12% is treated as an "
            "operational alert threshold for bounded bilingual deployments. Gaps above this threshold "
            "may indicate that the model performs materially differently across languages, warranting "
            "additional controls or human review."
        )

        doc.add_heading("Scope & Limitations", level=2)
        num_models = len(self.models)
        num_prompts = self.metadata.get("total_prompts", "N/A")
        doc.add_paragraph(
            f"This evaluation tested {num_models} model(s) against {num_prompts} prompts in a single "
            f"use case ({self.metadata.get('use_case', 'General')}). Results reflect observed performance "
            f"within this scope and may not generalize to other domains, prompt types, or model versions. "
            f"A broader evaluation is recommended before production deployment."
        )

        doc.add_page_break()

        # ═══════════════════════════════════════
        # DETAILED RESULTS
        # ═══════════════════════════════════════
        doc.add_heading("3. Detailed Results", level=1)

        # Get all unique dimensions from the data
        all_dims = set()
        for model_data in self.aggregates.values():
            all_dims.update(model_data.get("en", {}).keys())
            all_dims.update(model_data.get("ar", {}).keys())
        all_dims = sorted(all_dims)

        if all_dims and self.aggregates:
            doc.add_heading("3.1 Model Comparison", level=2)
            doc.add_paragraph(
                "The following table shows average scores per dimension for each model, "
                "comparing English and Arabic performance with the cross-lingual gap."
            )
            add_scores_table(doc, self.aggregates, all_dims)
            doc.add_paragraph()

        # ── Per-model analysis ──
        for idx, (model_id, model_data) in enumerate(self.aggregates.items()):
            model_scores = overall.get(model_id, {})
            short_name = (
                model_id.split("-")[0].upper() if "-" in model_id else model_id.upper()
            )

            doc.add_heading(f"3.{idx + 2} {short_name} Analysis", level=2)

            add_rich_paragraph(
                doc,
                [
                    {
                        "text": "Overall English Score: ",
                        "bold": True,
                        "color": GRAY,
                        "size": Pt(10),
                    },
                    {
                        "text": f"{model_scores.get('en_average', 'N/A')}%",
                        "bold": True,
                        "color": score_to_color(model_scores.get("en_average", 0)),
                        "size": Pt(11),
                    },
                    {
                        "text": "    Overall Arabic Score: ",
                        "bold": True,
                        "color": GRAY,
                        "size": Pt(10),
                    },
                    {
                        "text": f"{model_scores.get('ar_average', 'N/A')}%",
                        "bold": True,
                        "color": score_to_color(model_scores.get("ar_average", 0)),
                        "size": Pt(11),
                    },
                    {
                        "text": "    Average Gap: ",
                        "bold": True,
                        "color": GRAY,
                        "size": Pt(10),
                    },
                    {
                        "text": f"{model_scores.get('avg_gap', 'N/A')}%",
                        "bold": True,
                        "color": (
                            ORANGE if model_scores.get("avg_gap", 0) > 8 else GREEN
                        ),
                        "size": Pt(11),
                    },
                ],
            )

            # Collect flags for this model
            model_flags = self._collect_flags(model_id=model_id)

            if model_flags:
                doc.add_heading("Flagged Issues", level=3)

                # Group flags by dimension
                by_dimension = {}
                for f in model_flags:
                    dim = f["dimension"]
                    if dim not in by_dimension:
                        by_dimension[dim] = []
                    by_dimension[dim].append(f)

                for dim, dim_flags in by_dimension.items():
                    doc.add_heading(f"{dim.replace('_', ' ').title()} Flags", level=3)

                    # Show top 5 flags per dimension
                    for flag_data in dim_flags[:5]:
                        p = doc.add_paragraph(style="List Bullet")

                        lang_label = "AR" if flag_data["language"] == "ar" else "EN"
                        run = p.add_run(
                            f"[{lang_label}] [{flag_data['severity'].upper()}] "
                        )
                        run.font.size = Pt(9)
                        run.font.bold = True
                        run.font.color.rgb = severity_color(flag_data["severity"])

                        run = p.add_run(flag_data["flag"])
                        run.font.size = Pt(9)
                        run.font.color.rgb = DARK

                    if len(dim_flags) > 5:
                        add_styled_paragraph(
                            doc,
                            f"... and {len(dim_flags) - 5} additional {dim} flag(s)",
                            color=GRAY,
                            size=Pt(9),
                        )

        doc.add_page_break()

        # ═══════════════════════════════════════
        # CROSS-LINGUAL GAP ANALYSIS
        # ═══════════════════════════════════════
        doc.add_heading("4. Cross-Lingual Gap Analysis", level=1)

        doc.add_paragraph(
            "Cross-lingual gaps measure how differently each model performs between English and Arabic. "
            "Larger gaps indicate higher risk for bilingual deployments. Gaps above 12% are flagged as "
            "high-risk and may require language-specific guardrails."
        )

        for model_id, model_data in self.aggregates.items():
            short_name = (
                model_id.split("-")[0].upper() if "-" in model_id else model_id.upper()
            )
            gap_data = model_data.get("cross_lingual_gap", {})

            if gap_data:
                doc.add_heading(f"{short_name} Cross-Lingual Gaps", level=2)

                gap_rows = []
                for dim, gd in gap_data.items():
                    gap_rows.append(
                        (
                            dim.replace("_", " ").title(),
                            f"{gd.get('en_avg', 'N/A')}%",
                            f"{gd.get('ar_avg', 'N/A')}%",
                            f"{gd.get('gap', 'N/A')}%",
                            gd.get("severity", "N/A").title(),
                        )
                    )

                table = doc.add_table(rows=1 + len(gap_rows), cols=5)
                table.style = "Table Grid"

                for j, h in enumerate(
                    ["Dimension", "English", "Arabic", "Gap", "Severity"]
                ):
                    cell = table.cell(0, j)
                    run = cell.paragraphs[0].add_run(h)
                    run.font.name = FONT_BODY
                    run.font.size = Pt(9)
                    run.font.bold = True
                    run.font.color.rgb = WHITE
                    set_cell_shading(cell, "0C2D48")

                for i, row in enumerate(gap_rows):
                    for j, val in enumerate(row):
                        cell = table.cell(i + 1, j)
                        run = cell.paragraphs[0].add_run(str(val))
                        run.font.name = FONT_BODY
                        run.font.size = Pt(9)
                        if j == 4:  # Severity
                            run.font.bold = True
                            run.font.color.rgb = severity_color(val.lower())

                doc.add_paragraph()

        doc.add_page_break()

        # ═══════════════════════════════════════
        # RECOMMENDATIONS
        # ═══════════════════════════════════════
        doc.add_heading("5. Deployment Recommendation", level=1)

        # Get code-switching data for model-specific recommendations
        code_switching = self._get_code_switching_findings()
        cs_by_model = {}
        for cs in code_switching:
            if cs["model"] not in cs_by_model:
                cs_by_model[cs["model"]] = []
            cs_by_model[cs["model"]].append(cs)

        # Get bias flags per model
        bias_flags_by_model = {}
        for model_id in overall.keys():
            bias_flags_by_model[model_id] = self._collect_flags(
                model_id=model_id, dimension="bias"
            )

        # Recommended model callout
        recommended_model, reason = self._get_recommended_model(overall)
        if recommended_model and len(overall) > 1:
            rec_short = (
                recommended_model.split("-")[0].upper()
                if "-" in recommended_model
                else recommended_model.upper()
            )
            p = doc.add_paragraph()
            run = p.add_run(f"Recommended for this use case: {rec_short}")
            run.font.bold = True
            run.font.color.rgb = GREEN
            run.font.size = Pt(11)
            doc.add_paragraph(f"Rationale: {reason}")
            doc.add_paragraph()

        for model_id, scores in overall.items():
            short_name = (
                model_id.split("-")[0].upper() if "-" in model_id else model_id.upper()
            )
            status = scores["status"]
            critical_gaps = scores.get("critical_gaps", [])
            has_code_switching = model_id in cs_by_model
            bias_flags = bias_flags_by_model.get(model_id, [])

            # Build risk profile
            risk_factors = []
            if critical_gaps:
                risk_factors.append(
                    f"{len(critical_gaps)} critical cross-lingual gap(s)"
                )
            if has_code_switching:
                risk_factors.append(
                    f"{len(cs_by_model[model_id])} code-switching instance(s)"
                )
            if len(bias_flags) > 3:
                risk_factors.append(f"{len(bias_flags)} bias flag(s)")

            doc.add_heading(f"{short_name}: {status}", level=2)

            # Risk profile summary
            if risk_factors:
                p = doc.add_paragraph()
                run = p.add_run("Risk Profile: ")
                run.font.bold = True
                run.font.color.rgb = ORANGE if len(risk_factors) < 3 else RED
                run = p.add_run("; ".join(risk_factors))
                run.font.size = Pt(10)
            else:
                p = doc.add_paragraph()
                run = p.add_run("Risk Profile: ")
                run.font.bold = True
                run.font.color.rgb = GREEN
                run = p.add_run("Low risk — no critical issues detected")
                run.font.size = Pt(10)

            if status == "Ready for Pilot":
                # Differentiated messaging based on whether model has code-switching
                if has_code_switching:
                    doc.add_paragraph(
                        f"{short_name} demonstrates acceptable overall performance and is suitable for a bounded "
                        f"pilot with Arabic-language validation controls. While core metrics pass threshold, "
                        f"the detected code-switching ({len(cs_by_model.get(model_id, []))} instance(s)) requires "
                        f"a language validation layer before production deployment."
                    )
                else:
                    doc.add_paragraph(
                        f"{short_name} demonstrates acceptable performance across both languages for a bounded "
                        f"pilot deployment. Standard monitoring and periodic re-evaluation are recommended."
                    )
                doc.add_paragraph("Recommended controls:")
                controls = [
                    "Standard human review sampling (10-15% of outputs)",
                    "Monthly performance monitoring",
                    "Quarterly cross-lingual re-evaluation",
                    "User feedback collection mechanism",
                ]
                # Add specific controls for any detected issues
                if has_code_switching:
                    controls.insert(
                        0,
                        "Arabic language validation layer: verify responses are in requested language",
                    )
                    controls.insert(
                        1,
                        "Consider Arabic-specific system prompts to reduce language drift",
                    )
                for ctrl in controls:
                    doc.add_paragraph(ctrl, style="List Bullet")

            elif status == "Restricted Pilot Only":
                # Provide differentiated advice based on specific issues
                issues_text = []
                if critical_gaps:
                    gap_dims = [g[0] for g in critical_gaps]
                    issues_text.append(f"critical gaps in {', '.join(gap_dims)}")
                if has_code_switching:
                    issues_text.append("responds in English when prompted in Arabic")
                if len(bias_flags) > 3:
                    issues_text.append("multiple bias patterns detected")

                if issues_text:
                    doc.add_paragraph(
                        f"{short_name} requires additional controls due to: {'; '.join(issues_text)}. "
                        f"A restricted pilot with enhanced oversight is recommended."
                    )
                else:
                    doc.add_paragraph(
                        f"{short_name} shows moderate cross-lingual performance gaps that require additional "
                        f"controls before deployment. A restricted pilot with enhanced oversight is recommended."
                    )

                doc.add_paragraph("Required controls before deployment:")
                controls = []
                if critical_gaps:
                    controls.append(
                        "Enhanced human review for all Arabic outputs (100% review initially)"
                    )
                if has_code_switching:
                    controls.append(
                        "Language validation layer: reject or flag responses not in requested language"
                    )
                    controls.append(
                        "Consider Arabic-specific system prompts to reduce code-switching"
                    )
                if len(bias_flags) > 3:
                    controls.append(
                        "Domain-specific prompt engineering to reduce bias patterns"
                    )
                    controls.append(
                        "Implement bias-detection guardrails before responses reach users"
                    )
                controls.extend(
                    [
                        "Restricted scope: deploy in one workflow only before expanding",
                        "Weekly performance monitoring during pilot phase",
                        "Re-evaluation after 4 weeks with updated scoring",
                    ]
                )
                for ctrl in controls:
                    doc.add_paragraph(ctrl, style="List Bullet")

            else:
                specific_issues = []
                if critical_gaps:
                    for dim, gap in critical_gaps:
                        specific_issues.append(f"{dim} performance gap of {gap:.1f}%")
                if has_code_switching:
                    specific_issues.append(
                        "severe code-switching (responding in wrong language)"
                    )

                if specific_issues:
                    doc.add_paragraph(
                        f"{short_name} exhibits significant issues: {'; '.join(specific_issues)}. "
                        f"These make it unsuitable for bilingual deployment without substantial mitigation."
                    )
                else:
                    doc.add_paragraph(
                        f"{short_name} exhibits significant performance issues that make it unsuitable for "
                        f"bilingual deployment without substantial mitigation. Deployment is not recommended "
                        f"at this time."
                    )

                doc.add_paragraph("Recommended next steps:")
                next_steps = [
                    "Evaluate alternative model providers",
                    "Consider Arabic-specialist models (e.g. Jais, Allam) for Arabic-facing applications",
                ]
                if critical_gaps:
                    next_steps.append(
                        "Implement language-specific routing: use different models for Arabic vs English"
                    )
                if has_code_switching:
                    next_steps.append(
                        "Add explicit language enforcement in system prompts"
                    )
                next_steps.append("Re-evaluate after provider updates or fine-tuning")
                for ctrl in next_steps:
                    doc.add_paragraph(ctrl, style="List Bullet")

        doc.add_page_break()

        # ═══════════════════════════════════════
        # NEXT STEPS
        # ═══════════════════════════════════════
        doc.add_heading("6. Recommended Next Steps", level=1)

        steps = [
            (
                "Immediate (This Week)",
                [
                    "Review this report with your AI governance or technology leadership team",
                    "Identify the highest-risk findings and determine which require immediate action",
                    "Decide on deployment status for each evaluated model",
                ],
            ),
            (
                "Short-Term (Next 30 Days)",
                [
                    "Implement recommended controls for any models approved for pilot",
                    "Begin restricted pilot deployment if applicable",
                    "Establish monitoring and feedback collection processes",
                ],
            ),
            (
                "Medium-Term (60-90 Days)",
                [
                    "Conduct follow-up evaluation to measure improvement",
                    "Expand pilot scope if initial results are positive",
                    "Consider a Cross-Lingual Bias & Reliability Audit for deeper analysis",
                    "Evaluate additional models or providers as needed",
                ],
            ),
        ]

        for period, actions in steps:
            doc.add_heading(period, level=2)
            for action in actions:
                doc.add_paragraph(action, style="List Bullet")

        # ═══════════════════════════════════════
        # ABOUT / CONTACT
        # ═══════════════════════════════════════
        doc.add_page_break()
        doc.add_heading("About LinguaEval", level=1)

        doc.add_paragraph(
            "LinguaEval is a specialist multilingual AI evaluation and deployment studio based in the UK, "
            "serving organisations across the UK and GCC. We help organisations evaluate, de-risk, and deploy "
            "Arabic-English AI systems through research-grade bias auditing, cross-lingual benchmarking, "
            "and high-trust pilot delivery."
        )

        doc.add_paragraph()

        doc.add_heading("Our Services", level=2)
        services = [
            "Multilingual AI Readiness Assessment — Decision-ready evaluation with deployment recommendation",
            "Cross-Lingual Bias & Reliability Audit — Deep technical analysis with mitigation guidance",
            "High-Trust AI Pilot — Bounded deployment with governance and handover",
            "Assurance & Monitoring Retainer — Ongoing evaluation and advisory",
        ]
        for s in services:
            doc.add_paragraph(s, style="List Bullet")

        doc.add_paragraph()
        add_styled_paragraph(
            doc, "Contact: hello@linguaeval.com", color=ACCENT, bold=True
        )
        add_styled_paragraph(
            doc, "Website: www.linguaeval.com", color=ACCENT, bold=True
        )

        doc.add_paragraph()
        p = doc.add_paragraph()
        run = p.add_run(
            "\u00a9 2026 LinguaEval Ltd. All rights reserved. This report is confidential and prepared "
            "exclusively for the named client. It may not be distributed without written permission."
        )
        run.font.name = FONT_BODY
        run.font.size = Pt(9)
        run.font.color.rgb = GRAY
        run.font.italic = True

        # ── Save ──
        os.makedirs(
            os.path.dirname(output_path) if os.path.dirname(output_path) else ".",
            exist_ok=True,
        )
        doc.save(output_path)
        print(f"\n[+] Report generated: {output_path}")
        print(f"    Client: {self.client_name}")
        print(f"    Models: {', '.join(self.models)}")
        print(f"    Prompts: {self.metadata.get('total_prompts', 'N/A')}")

    def generate_pdf(self, output_path: str):
        """Generate executive summary PDF report."""
        overall = self._calculate_overall_scores()
        code_switching = self._get_code_switching_findings()
        recommended_model, reason = self._get_recommended_model(overall)

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()

        # ── Cover Page ──
        pdf.set_font("Helvetica", "B", 28)
        pdf.set_text_color(12, 45, 72)  # NAVY
        pdf.ln(40)
        pdf.cell(0, 15, "MULTILINGUAL AI", align="C", ln=True)
        pdf.cell(0, 15, "READINESS REPORT", align="C", ln=True)

        pdf.set_font("Helvetica", "", 12)
        pdf.set_text_color(100, 116, 139)  # GRAY
        use_case = self.metadata.get("use_case", "General")
        pdf.ln(5)
        pdf.cell(
            0, 8, f"Arabic-English Model Evaluation for {use_case}", align="C", ln=True
        )

        pdf.ln(10)
        pdf.set_font("Helvetica", "B", 14)
        pdf.set_text_color(2, 132, 199)  # ACCENT
        pdf.cell(0, 10, f"Prepared for {self.client_name}", align="C", ln=True)

        pdf.ln(15)
        pdf.set_draw_color(2, 132, 199)
        pdf.line(60, pdf.get_y(), 150, pdf.get_y())

        pdf.ln(10)
        pdf.set_font("Helvetica", "", 11)
        pdf.set_text_color(100, 116, 139)
        date_str = datetime.fromisoformat(self.timestamp).strftime("%B %Y")
        pdf.cell(0, 8, date_str, align="C", ln=True)
        pdf.cell(0, 8, "LinguaEval", align="C", ln=True)

        # ── Executive Summary Page ──
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 18)
        pdf.set_text_color(12, 45, 72)
        pdf.cell(0, 12, "1. Executive Summary", ln=True)
        pdf.ln(5)

        # Readiness status table
        pdf.set_font("Helvetica", "B", 10)
        pdf.set_fill_color(12, 45, 72)
        pdf.set_text_color(255, 255, 255)
        col_w = 63
        pdf.cell(col_w, 8, "Model", border=1, fill=True, align="C")
        pdf.cell(col_w, 8, "Readiness Status", border=1, fill=True, align="C")
        pdf.cell(col_w, 8, "Overall Score", border=1, fill=True, align="C")
        pdf.ln()

        pdf.set_font("Helvetica", "", 10)
        for model_id, scores in overall.items():
            short_name = (
                model_id.split("-")[0].upper() if "-" in model_id else model_id.upper()
            )
            status = scores["status"]
            avg_score = scores["combined_average"]

            pdf.set_text_color(30, 41, 59)
            pdf.cell(col_w, 8, short_name, border=1, align="C")

            # Status color
            if status == "Ready for Pilot":
                pdf.set_text_color(22, 101, 52)  # GREEN
            elif status == "Restricted Pilot Only":
                pdf.set_text_color(194, 65, 12)  # ORANGE
            else:
                pdf.set_text_color(185, 28, 28)  # RED
            pdf.cell(col_w, 8, status, border=1, align="C")

            pdf.set_text_color(30, 41, 59)
            pdf.cell(col_w, 8, f"{avg_score:.1f}%", border=1, align="C")
            pdf.ln()

        pdf.ln(8)

        # Recommended model
        if recommended_model and len(overall) > 1:
            rec_short = (
                recommended_model.split("-")[0].upper()
                if "-" in recommended_model
                else recommended_model.upper()
            )
            pdf.set_font("Helvetica", "B", 11)
            pdf.set_text_color(22, 101, 52)
            pdf.cell(0, 8, f"Recommended for this use case: {rec_short}", ln=True)
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(30, 41, 59)
            pdf.multi_cell(0, 6, f"Rationale: {reason}")
            pdf.ln(5)

        # Critical findings
        pdf.set_font("Helvetica", "B", 14)
        pdf.set_text_color(12, 45, 72)
        pdf.cell(0, 10, "Critical Findings", ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(30, 41, 59)

        if code_switching:
            models_with_cs = set(f["model"] for f in code_switching)
            pdf.set_text_color(185, 28, 28)
            pdf.set_font("Helvetica", "B", 10)
            pdf.cell(0, 6, "! Code-Switching Detected:", ln=True)
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(30, 41, 59)
            pdf.multi_cell(
                0,
                5,
                f"{len(code_switching)} instance(s) where Arabic prompts received responses with significant English text. Affected: {', '.join(models_with_cs)}.",
            )
            pdf.ln(3)

        critical_findings = []
        for model_id, scores in overall.items():
            for dim, gap in scores.get("critical_gaps", []):
                short = (
                    model_id.split("-")[0].upper()
                    if "-" in model_id
                    else model_id.upper()
                )
                critical_findings.append(f"{short}: {dim} gap of {gap:.1f}%")

        if critical_findings:
            pdf.set_text_color(185, 28, 28)
            pdf.set_font("Helvetica", "B", 10)
            pdf.cell(0, 6, "! Critical Cross-Lingual Gaps:", ln=True)
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(30, 41, 59)
            pdf.multi_cell(0, 5, "; ".join(critical_findings))

        if not code_switching and not critical_findings:
            pdf.set_text_color(22, 101, 52)
            pdf.cell(0, 6, "No critical issues detected.", ln=True)

        # ── Methodology Page ──
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 18)
        pdf.set_text_color(12, 45, 72)
        pdf.cell(0, 12, "2. Methodology", ln=True)
        pdf.ln(3)

        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(30, 41, 59)
        pdf.multi_cell(
            0,
            5,
            "This assessment used the LinguaEval evaluation framework, which tests AI models across six dimensions in both Arabic and English using sector-specific prompt packs.",
        )
        pdf.ln(5)

        # Parameters table
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 8, "Evaluation Parameters", ln=True)
        pdf.set_font("Helvetica", "", 10)

        params = [
            ("Client", self.client_name),
            ("Sector", self.sector.title()),
            ("Use Case", self.metadata.get("use_case", "General")),
            ("Models Evaluated", ", ".join(self.models)),
            ("Languages", "English, Arabic"),
            ("Total Prompts", str(self.metadata.get("total_prompts", "N/A"))),
            (
                "Evaluation Date",
                datetime.fromisoformat(self.timestamp).strftime("%d %B %Y"),
            ),
        ]

        for key, val in params:
            pdf.set_font("Helvetica", "B", 10)
            pdf.cell(50, 6, key + ":", border=0)
            pdf.set_font("Helvetica", "", 10)
            pdf.cell(0, 6, str(val), ln=True)

        pdf.ln(8)

        # Scoring threshold note
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 8, "Scoring Thresholds", ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(
            0,
            5,
            "Within the LinguaEval framework, a cross-lingual gap above 12% is treated as an operational alert threshold for bounded bilingual deployments.",
        )

        pdf.ln(8)

        # Scope note
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 8, "Scope & Limitations", ln=True)
        pdf.set_font("Helvetica", "", 10)
        num_models = len(self.models)
        num_prompts = self.metadata.get("total_prompts", "N/A")
        pdf.multi_cell(
            0,
            5,
            f"This evaluation tested {num_models} model(s) against {num_prompts} prompts in a single use case. Results reflect observed performance within this scope and may not generalize to other domains.",
        )

        # ── Results Summary Page ──
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 18)
        pdf.set_text_color(12, 45, 72)
        pdf.cell(0, 12, "3. Results Summary", ln=True)
        pdf.ln(3)

        # Get all dimensions
        all_dims = set()
        for model_data in self.aggregates.values():
            all_dims.update(model_data.get("en", {}).keys())
            all_dims.update(model_data.get("ar", {}).keys())
        all_dims = sorted(all_dims)

        # Score table per model
        for model_id, model_data in self.aggregates.items():
            short_name = (
                model_id.split("-")[0].upper() if "-" in model_id else model_id.upper()
            )
            pdf.set_font("Helvetica", "B", 12)
            pdf.set_text_color(12, 45, 72)
            pdf.cell(0, 8, f"{short_name} Scores by Dimension", ln=True)

            # Table header
            pdf.set_font("Helvetica", "B", 9)
            pdf.set_fill_color(12, 45, 72)
            pdf.set_text_color(255, 255, 255)
            pdf.cell(50, 7, "Dimension", border=1, fill=True)
            pdf.cell(30, 7, "English", border=1, fill=True, align="C")
            pdf.cell(30, 7, "Arabic", border=1, fill=True, align="C")
            pdf.cell(30, 7, "Gap", border=1, fill=True, align="C")
            pdf.ln()

            pdf.set_font("Helvetica", "", 9)
            for dim in all_dims:
                en_score = model_data.get("en", {}).get(dim, {}).get("average", "N/A")
                ar_score = model_data.get("ar", {}).get(dim, {}).get("average", "N/A")
                gap_data = model_data.get("cross_lingual_gap", {}).get(dim, {})
                gap = gap_data.get("gap", "N/A")

                pdf.set_text_color(30, 41, 59)
                pdf.cell(50, 6, dim.replace("_", " ").title(), border=1)

                # EN score
                if isinstance(en_score, (int, float)):
                    pdf.cell(30, 6, f"{en_score:.1f}%", border=1, align="C")
                else:
                    pdf.cell(30, 6, "N/A", border=1, align="C")

                # AR score
                if isinstance(ar_score, (int, float)):
                    pdf.cell(30, 6, f"{ar_score:.1f}%", border=1, align="C")
                else:
                    pdf.cell(30, 6, "N/A", border=1, align="C")

                # Gap with color
                if isinstance(gap, (int, float)):
                    if gap > 12:
                        pdf.set_text_color(185, 28, 28)  # RED
                    elif gap > 8:
                        pdf.set_text_color(194, 65, 12)  # ORANGE
                    else:
                        pdf.set_text_color(22, 101, 52)  # GREEN
                    pdf.cell(30, 6, f"{gap:.1f}%", border=1, align="C")
                else:
                    pdf.set_text_color(100, 116, 139)
                    pdf.cell(30, 6, "-", border=1, align="C")
                pdf.ln()

            pdf.ln(8)

        # ── Recommendations Page ──
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 18)
        pdf.set_text_color(12, 45, 72)
        pdf.cell(0, 12, "4. Deployment Recommendations", ln=True)
        pdf.ln(3)

        for model_id, scores in overall.items():
            short_name = (
                model_id.split("-")[0].upper() if "-" in model_id else model_id.upper()
            )
            status = scores["status"]

            pdf.set_font("Helvetica", "B", 12)
            if status == "Ready for Pilot":
                pdf.set_text_color(22, 101, 52)
            elif status == "Restricted Pilot Only":
                pdf.set_text_color(194, 65, 12)
            else:
                pdf.set_text_color(185, 28, 28)
            pdf.cell(0, 8, f"{short_name}: {status}", ln=True)

            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(30, 41, 59)

            if status == "Ready for Pilot":
                pdf.multi_cell(
                    0,
                    5,
                    f"{short_name} demonstrates acceptable performance across both languages for a bounded pilot deployment. Standard monitoring and periodic re-evaluation are recommended.",
                )
            elif status == "Restricted Pilot Only":
                pdf.multi_cell(
                    0,
                    5,
                    f"{short_name} shows moderate cross-lingual performance gaps requiring additional controls before deployment. A restricted pilot with enhanced oversight is recommended.",
                )
            else:
                pdf.multi_cell(
                    0,
                    5,
                    f"{short_name} exhibits significant performance issues that make it unsuitable for bilingual deployment without substantial mitigation.",
                )

            pdf.ln(5)

        # ── Footer ──
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 14)
        pdf.set_text_color(12, 45, 72)
        pdf.cell(0, 10, "About LinguaEval", ln=True)

        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(30, 41, 59)
        pdf.multi_cell(
            0,
            5,
            "LinguaEval is a specialist multilingual AI evaluation and deployment studio based in the UK, serving organisations across the UK and GCC. We help organisations evaluate, de-risk, and deploy Arabic-English AI systems.",
        )

        pdf.ln(8)
        pdf.set_font("Helvetica", "B", 10)
        pdf.set_text_color(2, 132, 199)
        pdf.cell(0, 6, "Contact: hello@linguaeval.com", ln=True)
        pdf.cell(0, 6, "Website: www.linguaeval.com", ln=True)

        pdf.ln(10)
        pdf.set_font("Helvetica", "I", 8)
        pdf.set_text_color(100, 116, 139)
        pdf.multi_cell(
            0,
            4,
            "2026 LinguaEval Ltd. All rights reserved. This report is confidential and prepared exclusively for the named client.",
        )

        # ── Save PDF ──
        os.makedirs(
            os.path.dirname(output_path) if os.path.dirname(output_path) else ".",
            exist_ok=True,
        )
        pdf.output(output_path)
        print(f"\n[+] PDF Report generated: {output_path}")
        print(f"    Client: {self.client_name}")
        print(f"    Models: {', '.join(self.models)}")
        print(f"    Prompts: {self.metadata.get('total_prompts', 'N/A')}")

    def generate_slides(self, output_path: str):
        """Generate executive summary slide deck (PPTX)."""
        overall = self._calculate_overall_scores()
        code_switching = self._get_code_switching_findings()
        recommended_model, reason = self._get_recommended_model(overall)

        prs = Presentation()
        prs.slide_width = PptxInches(13.333)  # 16:9 widescreen
        prs.slide_height = PptxInches(7.5)

        # Color constants for slides
        PPTX_NAVY = PptxRGB(0x0C, 0x2D, 0x48)
        PPTX_ACCENT = PptxRGB(0x02, 0x84, 0xC7)
        PPTX_GREEN = PptxRGB(0x16, 0x65, 0x34)
        PPTX_RED = PptxRGB(0xB9, 0x1C, 0x1C)
        PPTX_ORANGE = PptxRGB(0xC2, 0x41, 0x0C)
        PPTX_GRAY = PptxRGB(0x64, 0x74, 0x8B)
        PPTX_WHITE = PptxRGB(0xFF, 0xFF, 0xFF)

        # ══════════════════════════════════════════
        # SLIDE 1: Title Slide
        # ══════════════════════════════════════════
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(2.5), PptxInches(12.333), PptxInches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "MULTILINGUAL AI READINESS REPORT"
        p.font.size = PptxPt(44)
        p.font.bold = True
        p.font.color.rgb = PPTX_NAVY
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        use_case = self.metadata.get("use_case", "General")
        sub_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(3.6), PptxInches(12.333), PptxInches(0.5)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Arabic-English Model Evaluation for {use_case}"
        p.font.size = PptxPt(20)
        p.font.color.rgb = PPTX_GRAY
        p.alignment = PP_ALIGN.CENTER

        # Client name
        client_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(4.3), PptxInches(12.333), PptxInches(0.5)
        )
        tf = client_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Prepared for {self.client_name}"
        p.font.size = PptxPt(18)
        p.font.bold = True
        p.font.color.rgb = PPTX_ACCENT
        p.alignment = PP_ALIGN.CENTER

        # Date
        date_str = datetime.fromisoformat(self.timestamp).strftime("%B %Y")
        date_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(5.5), PptxInches(12.333), PptxInches(0.4)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{date_str} | LinguaEval"
        p.font.size = PptxPt(14)
        p.font.color.rgb = PPTX_GRAY
        p.alignment = PP_ALIGN.CENTER

        # ══════════════════════════════════════════
        # SLIDE 2: Executive Summary
        # ══════════════════════════════════════════
        slide = prs.slides.add_slide(slide_layout)

        # Header
        header_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.4), PptxInches(12.333), PptxInches(0.6)
        )
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = PptxPt(32)
        p.font.bold = True
        p.font.color.rgb = PPTX_NAVY

        # Summary table - create a table for model status
        table = slide.shapes.add_table(
            rows=1 + len(overall),
            cols=3,
            left=PptxInches(1),
            top=PptxInches(1.4),
            width=PptxInches(11),
            height=PptxInches(0.5 + 0.5 * len(overall)),
        ).table

        # Header row
        for i, header in enumerate(["Model", "Readiness Status", "Overall Score"]):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = PPTX_NAVY
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = PptxPt(14)
            p.font.color.rgb = PPTX_WHITE
            p.alignment = PP_ALIGN.CENTER

        # Data rows
        for row_idx, (model_id, scores) in enumerate(overall.items(), start=1):
            short_name = (
                model_id.split("-")[0].upper() if "-" in model_id else model_id.upper()
            )
            status = scores["status"]
            avg_score = scores["combined_average"]

            # Model name
            cell = table.cell(row_idx, 0)
            cell.text = short_name
            p = cell.text_frame.paragraphs[0]
            p.font.size = PptxPt(12)
            p.alignment = PP_ALIGN.CENTER

            # Status
            cell = table.cell(row_idx, 1)
            cell.text = status
            p = cell.text_frame.paragraphs[0]
            p.font.size = PptxPt(12)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            if status == "Ready for Pilot":
                p.font.color.rgb = PPTX_GREEN
            elif status == "Restricted Pilot Only":
                p.font.color.rgb = PPTX_ORANGE
            else:
                p.font.color.rgb = PPTX_RED

            # Score
            cell = table.cell(row_idx, 2)
            cell.text = f"{avg_score:.1f}%"
            p = cell.text_frame.paragraphs[0]
            p.font.size = PptxPt(12)
            p.alignment = PP_ALIGN.CENTER

        # Recommended model callout
        if recommended_model and len(overall) > 1:
            rec_short = (
                recommended_model.split("-")[0].upper()
                if "-" in recommended_model
                else recommended_model.upper()
            )
            rec_box = slide.shapes.add_textbox(
                PptxInches(1),
                PptxInches(3.2 + 0.5 * len(overall)),
                PptxInches(11),
                PptxInches(0.5),
            )
            tf = rec_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"✓ Recommended: {rec_short}"
            p.font.size = PptxPt(16)
            p.font.bold = True
            p.font.color.rgb = PPTX_GREEN

        # Critical findings
        findings_y = 4.0 + 0.5 * len(overall)
        if code_switching or any(
            scores.get("critical_gaps") for scores in overall.values()
        ):
            findings_box = slide.shapes.add_textbox(
                PptxInches(1), PptxInches(findings_y), PptxInches(11), PptxInches(1.5)
            )
            tf = findings_box.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = "Critical Findings:"
            p.font.size = PptxPt(14)
            p.font.bold = True
            p.font.color.rgb = PPTX_RED

            if code_switching:
                p = tf.add_paragraph()
                models_with_cs = set(f["model"] for f in code_switching)
                p.text = f"• Code-switching: {len(code_switching)} instance(s) in {', '.join(models_with_cs)}"
                p.font.size = PptxPt(12)

            for model_id, scores in overall.items():
                for dim, gap in scores.get("critical_gaps", []):
                    p = tf.add_paragraph()
                    short = (
                        model_id.split("-")[0].upper()
                        if "-" in model_id
                        else model_id.upper()
                    )
                    p.text = f"• {short}: {dim} cross-lingual gap of {gap:.1f}%"
                    p.font.size = PptxPt(12)

        # ══════════════════════════════════════════
        # SLIDE 3: Model Comparison
        # ══════════════════════════════════════════
        slide = prs.slides.add_slide(slide_layout)

        header_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.4), PptxInches(12.333), PptxInches(0.6)
        )
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Model Performance Comparison"
        p.font.size = PptxPt(32)
        p.font.bold = True
        p.font.color.rgb = PPTX_NAVY

        # Get all dimensions
        all_dims = set()
        for model_data in self.aggregates.values():
            all_dims.update(model_data.get("en", {}).keys())
            all_dims.update(model_data.get("ar", {}).keys())
        all_dims = sorted(all_dims)

        # Create comparison table
        models = list(self.aggregates.keys())
        table = slide.shapes.add_table(
            rows=1 + len(all_dims),
            cols=1 + len(models) * 2,
            left=PptxInches(0.5),
            top=PptxInches(1.2),
            width=PptxInches(12.333),
            height=PptxInches(0.4 + 0.4 * len(all_dims)),
        ).table

        # Header row
        cell = table.cell(0, 0)
        cell.text = "Dimension"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PPTX_NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = PptxPt(11)
        p.font.color.rgb = PPTX_WHITE

        col_idx = 1
        for model_id in models:
            short = (
                model_id.split("-")[0].upper() if "-" in model_id else model_id.upper()
            )
            for lang in ["EN", "AR"]:
                cell = table.cell(0, col_idx)
                cell.text = f"{short} {lang}"
                cell.fill.solid()
                cell.fill.fore_color.rgb = PPTX_NAVY
                p = cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = PptxPt(10)
                p.font.color.rgb = PPTX_WHITE
                p.alignment = PP_ALIGN.CENTER
                col_idx += 1

        # Data rows
        for row_idx, dim in enumerate(all_dims, start=1):
            cell = table.cell(row_idx, 0)
            cell.text = dim.replace("_", " ").title()
            p = cell.text_frame.paragraphs[0]
            p.font.size = PptxPt(10)
            p.font.bold = True

            col_idx = 1
            for model_id in models:
                model_data = self.aggregates.get(model_id, {})
                en_score = model_data.get("en", {}).get(dim, {}).get("average", None)
                ar_score = model_data.get("ar", {}).get(dim, {}).get("average", None)

                for score in [en_score, ar_score]:
                    cell = table.cell(row_idx, col_idx)
                    if score is not None:
                        cell.text = f"{score:.0f}%"
                        p = cell.text_frame.paragraphs[0]
                        p.font.size = PptxPt(10)
                        p.alignment = PP_ALIGN.CENTER
                        if score >= 80:
                            p.font.color.rgb = PPTX_GREEN
                        elif score >= 65:
                            p.font.color.rgb = PPTX_ORANGE
                        else:
                            p.font.color.rgb = PPTX_RED
                    else:
                        cell.text = "—"
                        p = cell.text_frame.paragraphs[0]
                        p.font.size = PptxPt(10)
                        p.alignment = PP_ALIGN.CENTER
                        p.font.color.rgb = PPTX_GRAY
                    col_idx += 1

        # ══════════════════════════════════════════
        # SLIDE 4: Recommendations
        # ══════════════════════════════════════════
        slide = prs.slides.add_slide(slide_layout)

        header_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.4), PptxInches(12.333), PptxInches(0.6)
        )
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Deployment Recommendations"
        p.font.size = PptxPt(32)
        p.font.bold = True
        p.font.color.rgb = PPTX_NAVY

        y_pos = 1.3
        for model_id, scores in overall.items():
            short_name = (
                model_id.split("-")[0].upper() if "-" in model_id else model_id.upper()
            )
            status = scores["status"]

            # Model header
            model_box = slide.shapes.add_textbox(
                PptxInches(0.7), PptxInches(y_pos), PptxInches(11.5), PptxInches(0.4)
            )
            tf = model_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{short_name}: {status}"
            p.font.size = PptxPt(18)
            p.font.bold = True
            if status == "Ready for Pilot":
                p.font.color.rgb = PPTX_GREEN
            elif status == "Restricted Pilot Only":
                p.font.color.rgb = PPTX_ORANGE
            else:
                p.font.color.rgb = PPTX_RED

            y_pos += 0.5

            # Recommendation text
            rec_box = slide.shapes.add_textbox(
                PptxInches(0.9), PptxInches(y_pos), PptxInches(11.3), PptxInches(1)
            )
            tf = rec_box.text_frame
            tf.word_wrap = True

            if status == "Ready for Pilot":
                p = tf.paragraphs[0]
                p.text = "• Suitable for bounded pilot deployment"
                p.font.size = PptxPt(12)
                p = tf.add_paragraph()
                p.text = "• Standard monitoring recommended (10-15% human review)"
                p.font.size = PptxPt(12)
                p = tf.add_paragraph()
                p.text = "• Quarterly cross-lingual re-evaluation"
                p.font.size = PptxPt(12)
            elif status == "Restricted Pilot Only":
                p = tf.paragraphs[0]
                p.text = "• Enhanced human review required initially"
                p.font.size = PptxPt(12)
                p = tf.add_paragraph()
                p.text = "• Restrict to single workflow before expanding"
                p.font.size = PptxPt(12)
                p = tf.add_paragraph()
                p.text = "• Weekly monitoring during pilot phase"
                p.font.size = PptxPt(12)
            else:
                p = tf.paragraphs[0]
                p.text = "• Not recommended for deployment at this time"
                p.font.size = PptxPt(12)
                p = tf.add_paragraph()
                p.text = "• Consider alternative models or fine-tuning"
                p.font.size = PptxPt(12)
                p = tf.add_paragraph()
                p.text = "• Re-evaluate after mitigation steps"
                p.font.size = PptxPt(12)

            y_pos += 1.3

        # ══════════════════════════════════════════
        # SLIDE 5: Next Steps / Contact
        # ══════════════════════════════════════════
        slide = prs.slides.add_slide(slide_layout)

        header_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.4), PptxInches(12.333), PptxInches(0.6)
        )
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Next Steps"
        p.font.size = PptxPt(32)
        p.font.bold = True
        p.font.color.rgb = PPTX_NAVY

        steps_box = slide.shapes.add_textbox(
            PptxInches(0.7), PptxInches(1.3), PptxInches(11.5), PptxInches(3)
        )
        tf = steps_box.text_frame
        tf.word_wrap = True

        steps = [
            (
                "Immediate (0-30 Days)",
                [
                    "Review critical findings with stakeholders",
                    "Implement recommended controls for pilot",
                    "Establish monitoring and feedback processes",
                ],
            ),
            (
                "Medium-Term (60-90 Days)",
                [
                    "Conduct follow-up evaluation",
                    "Expand pilot scope if results positive",
                    "Consider Cross-Lingual Bias & Reliability Audit",
                ],
            ),
        ]

        for period, actions in steps:
            p = tf.add_paragraph()
            p.text = period
            p.font.size = PptxPt(14)
            p.font.bold = True
            p.font.color.rgb = PPTX_ACCENT

            for action in actions:
                p = tf.add_paragraph()
                p.text = f"  • {action}"
                p.font.size = PptxPt(12)

        # Contact info
        contact_box = slide.shapes.add_textbox(
            PptxInches(0.7), PptxInches(5.5), PptxInches(11.5), PptxInches(1)
        )
        tf = contact_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Contact: hello@linguaeval.com | www.linguaeval.com"
        p.font.size = PptxPt(14)
        p.font.bold = True
        p.font.color.rgb = PPTX_ACCENT
        p.alignment = PP_ALIGN.CENTER

        # ── Save PPTX ──
        os.makedirs(
            os.path.dirname(output_path) if os.path.dirname(output_path) else ".",
            exist_ok=True,
        )
        prs.save(output_path)
        print(f"\n[+] Slide deck generated: {output_path}")
        print(f"    Client: {self.client_name}")
        print(f"    Models: {', '.join(self.models)}")
        print(f"    Slides: 5")


def main():
    parser = argparse.ArgumentParser(description="LinguaEval Report Generator")
    parser.add_argument(
        "--results", required=True, help="Path to evaluation results JSON"
    )
    parser.add_argument("--output", default=None, help="Output path for report")
    parser.add_argument(
        "--format",
        choices=["docx", "pdf", "slides", "both", "all"],
        default="docx",
        help="Output format: docx, pdf, slides, both (docx+pdf), or all (default: docx)",
    )

    args = parser.parse_args()

    if not os.path.exists(args.results):
        print(f"[!] Results file not found: {args.results}")
        sys.exit(1)

    base = os.path.splitext(os.path.basename(args.results))[0]

    generator = ReportGenerator(args.results)

    if args.format in ("docx", "both", "all"):
        docx_output = (
            args.output
            if args.output and args.output.endswith(".docx")
            else os.path.join("reports", f"{base}_report.docx")
        )
        generator.generate(docx_output)

    if args.format in ("pdf", "both", "all"):
        pdf_output = (
            args.output
            if args.output and args.output.endswith(".pdf")
            else os.path.join("reports", f"{base}_report.pdf")
        )
        generator.generate_pdf(pdf_output)

    if args.format in ("slides", "all"):
        slides_output = (
            args.output
            if args.output and args.output.endswith(".pptx")
            else os.path.join("reports", f"{base}_report.pptx")
        )
        generator.generate_slides(slides_output)


if __name__ == "__main__":
    main()
